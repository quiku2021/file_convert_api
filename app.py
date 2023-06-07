from flask import Flask, request
from flask_cors import CORS
import os
import PyPDF2
from docx import Document
from pptx import Presentation
from pyth.plugins.plaintext.writer import PlaintextWriter
from pyth.plugins.rtf15.reader import Rtf15Reader

app = Flask(__name__)
CORS(app)

@app.route("/")
def home():
    return "Bienvenidos a File Convert API "

@app.route('/convert', methods=['POST'])
def convert_to_text():
    file = request.files['file']
    if file:
        if file.filename.endswith('.pdf'):
            text = convert_pdf_to_text(file)
        elif file.filename.endswith('.docx'):
            text = convert_docx_to_text(file)
        elif file.filename.endswith('.pptx'):
            text = convert_pptx_to_text(file)
        elif file.filename.endswith('.rtf'):
            text = convert_rtf_to_text(file)
        else:
            return "Tipo de archivo no válido"
        return text
    else:
        return "Archivo no válido"

def convert_pdf_to_text(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def convert_docx_to_text(file):
    document = Document(file)
    text = '\n'.join([paragraph.text for paragraph in document.paragraphs])
    return text

def convert_pptx_to_text(file):
    presentation = Presentation(file)
    textos = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        textos.append(text)
    return textos

def convert_rtf_to_text(file):
    url_temp = './temp/' + file.filename
    file.save(url_temp)
    documentOpen = open(url_temp, 'rb')
    documentRead = Rtf15Reader.read(documentOpen)
    text = PlaintextWriter.write(documentRead).getvalue()
    documentOpen.close()
    os.remove(url_temp)   
    return text

if __name__ == '__main__':
    app.run()